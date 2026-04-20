"""Bank-style one-pager PPTX — sell-side equity research teaser layout.

Single landscape slide (16:9) with:
  • Header: Topaz wordmark (left) + ticker/company (center) + recommendation chip (right)
  • KPI strip: 5 big-number financial highlights
  • Left column: Investment Thesis + Key Strengths + Key Risks
  • Center column: Valuation Snapshot (DCF Bear/Base/Bull + analyst target)
  • Right column: Football field miniature + Revenue/consensus chart
  • Footer: data source + generation date + confidentiality notice
"""
from __future__ import annotations
from datetime import date
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ..extract.income_statement import IncomeStatement
from ..data.yfinance_client import QuoteSnapshot
from .chart import revenue_kpi_chart


# ---------- Brand palette ----------
# Navy-forward, typical of sell-side research (MS, GS, JPM)
COLOR_NAVY   = RGBColor(0x1F, 0x38, 0x64)     # primary / header bg
COLOR_NAVY_D = RGBColor(0x14, 0x24, 0x40)     # deeper navy
COLOR_ACCENT = RGBColor(0xC9, 0xA2, 0x27)     # muted gold accent (Topaz theme)
COLOR_GREEN  = RGBColor(0x2E, 0x7D, 0x32)     # BUY color
COLOR_ORANGE = RGBColor(0xE0, 0x7B, 0x00)     # HOLD / WAIT color
COLOR_RED    = RGBColor(0xC0, 0x39, 0x2B)     # SELL / AVOID color
COLOR_LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)     # near-white row bg
COLOR_GRAY   = RGBColor(0x59, 0x59, 0x59)
COLOR_DARK_G = RGBColor(0x37, 0x41, 0x51)
COLOR_BLACK  = RGBColor(0x11, 0x18, 0x22)
COLOR_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


# ==================================================================
# Low-level helpers
# ==================================================================

def _add_text(slide, left, top, width, height, text, *, size=10, bold=False,
              color=COLOR_BLACK, align=PP_ALIGN.LEFT, italic=False,
              font_name="Calibri", vertical=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = vertical
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def _add_rect(slide, left, top, width, height, fill, line_color=None, line_width=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        if line_width:
            s.line.width = line_width
    s.shadow.inherit = False
    return s


def _add_rounded_rect(slide, left, top, width, height, fill, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
    s.shadow.inherit = False
    return s


def _add_kpi_cell(slide, left, top, width, height, label, value):
    """A KPI box: small uppercase label + large value."""
    _add_rect(slide, left, top, width, height, COLOR_LIGHT)
    # Thin navy strip on left edge
    _add_rect(slide, left, top, Inches(0.04), height, COLOR_NAVY)
    # Label
    _add_text(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.2),
              Inches(0.22), label.upper(), size=8, bold=True, color=COLOR_GRAY)
    # Value
    _add_text(slide, left + Inches(0.15), top + Inches(0.32), width - Inches(0.2),
              Inches(0.45), value, size=16, bold=True, color=COLOR_NAVY_D)


def _add_section_header(slide, left, top, width, text, color=COLOR_NAVY):
    """A small section title with a thin underline bar."""
    _add_text(slide, left, top, width, Inches(0.28), text.upper(),
              size=10, bold=True, color=color)
    _add_rect(slide, left, top + Inches(0.27), Inches(0.4), Emu(18000), COLOR_ACCENT)


def _add_bullet_block(slide, left, top, width, height, bullets: List[str],
                      text_color=COLOR_BLACK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = Emu(0)
    tf.margin_right = tf.margin_bottom = Emu(0)
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = f"▪  {b}"
        run.font.size = Pt(9)
        run.font.color.rgb = text_color


# ==================================================================
# Recommendation chip
# ==================================================================

def _recommendation_chip(slide, left, top, width, height, recommendation: str,
                         target_price: Optional[float], upside: Optional[float]):
    """Top-right recommendation pill with target price + upside."""
    rec = (recommendation or "").upper()
    if any(k in rec for k in ("INVEST NOW", "BUY", "STRONG BUY")):
        label, color = "BUY", COLOR_GREEN
    elif any(k in rec for k in ("WAIT", "HOLD", "NEUTRAL")):
        label, color = "HOLD", COLOR_ORANGE
    elif any(k in rec for k in ("DON", "DO NOT", "AVOID", "SELL")):
        label, color = "AVOID", COLOR_RED
    else:
        label, color = "REVIEW", COLOR_DARK_G

    # Colored chip
    _add_rounded_rect(slide, left, top, width, height, color)

    # Top line: big BUY/HOLD/AVOID text
    _add_text(slide, left, top + Inches(0.08), width, Inches(0.32),
              label, size=16, bold=True, color=COLOR_WHITE,
              align=PP_ALIGN.CENTER)
    # Second line: price target + upside
    if target_price is not None:
        tgt_text = f"Target: ${target_price:.2f}"
        if upside is not None:
            tgt_text += f"   ({upside:+.0%})"
        _add_text(slide, left, top + Inches(0.40), width, Inches(0.25),
                  tgt_text, size=9, bold=True, color=COLOR_WHITE,
                  align=PP_ALIGN.CENTER)


# ==================================================================
# Valuation snapshot table
# ==================================================================

def _valuation_snapshot(slide, left, top, width, height, *,
                        current_price, dcf_bear, dcf_base, dcf_bull,
                        analyst_low, analyst_mean, analyst_high, analyst_count):
    """Compact price target summary: DCF scenarios + analyst consensus."""
    # Outer container
    _add_rect(slide, left, top, width, height, COLOR_WHITE,
              line_color=COLOR_NAVY)
    # Header bar
    _add_rect(slide, left, top, width, Inches(0.32), COLOR_NAVY)
    _add_text(slide, left + Inches(0.12), top + Inches(0.05),
              width - Inches(0.24), Inches(0.25),
              "VALUATION SNAPSHOT", size=10, bold=True, color=COLOR_WHITE)

    # Rows — 2 columns per row: label + value
    rows = []
    if current_price is not None:
        rows.append(("Current price", f"${current_price:.2f}", None, True))
    rows.append(("", "", None, False))  # spacer
    rows.append(("DCF — Bear case", _fmt_px(dcf_bear), _upside(dcf_bear, current_price), False))
    rows.append(("DCF — Base case", _fmt_px(dcf_base), _upside(dcf_base, current_price), True))
    rows.append(("DCF — Bull case", _fmt_px(dcf_bull), _upside(dcf_bull, current_price), False))
    rows.append(("", "", None, False))  # spacer
    rows.append((f"Analyst low ({analyst_count} analysts)" if analyst_count else "Analyst low",
                 _fmt_px(analyst_low), _upside(analyst_low, current_price), False))
    rows.append(("Analyst mean target", _fmt_px(analyst_mean),
                 _upside(analyst_mean, current_price), True))
    rows.append(("Analyst high", _fmt_px(analyst_high),
                 _upside(analyst_high, current_price), False))

    row_h = (height - Inches(0.35)) / max(len(rows), 1)
    y = top + Inches(0.36)
    for label, value, upside_str, bold in rows:
        if not label and not value:
            y += row_h
            continue
        _add_text(slide, left + Inches(0.12), y, width * 0.55, row_h,
                  label, size=9, bold=bold, color=COLOR_DARK_G,
                  vertical=MSO_ANCHOR.MIDDLE)
        _add_text(slide, left + width * 0.55, y, width * 0.25, row_h,
                  value, size=10, bold=bold, color=COLOR_NAVY_D,
                  align=PP_ALIGN.RIGHT, vertical=MSO_ANCHOR.MIDDLE)
        if upside_str:
            _add_text(slide, left + width * 0.80, y, width * 0.18, row_h,
                      upside_str, size=9, bold=False,
                      color=COLOR_GREEN if (upside_str.startswith("+") and upside_str != "+0%") else
                            COLOR_RED if upside_str.startswith("-") else COLOR_DARK_G,
                      align=PP_ALIGN.RIGHT, vertical=MSO_ANCHOR.MIDDLE)
        y += row_h


def _fmt_px(v: Optional[float]) -> str:
    if v is None:
        return "—"
    return f"${v:,.2f}"


def _upside(target: Optional[float], current: Optional[float]) -> Optional[str]:
    if target is None or current is None or current == 0:
        return None
    u = target / current - 1
    return f"{u:+.0%}"


# ==================================================================
# Revenue series extraction (unchanged)
# ==================================================================

def _quarterly_revenue_series(stmt: IncomeStatement) -> Dict[str, float]:
    result = {}
    rev_values = {}
    for line in stmt.lines:
        if line.section == "revenue_total":
            rev_values = line.values
            break
    for p in stmt.periods:
        if "Q" in p.label and p.label in rev_values:
            result[p.label] = rev_values[p.label]
    return result


# ==================================================================
# Main build
# ==================================================================

def build_onepager(stmt: IncomeStatement, quote: QuoteSnapshot,
                   one_pager_content: Optional[Dict] = None,
                   output_path: Path = None,
                   dcf=None,
                   football_field_png: Optional[Path] = None) -> Path:
    """Build a single landscape PPTX slide — bank-style equity research teaser."""
    one_pager_content = one_pager_content or {}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ===== Background tint =====
    _add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLOR_WHITE)

    # ===== HEADER (0 – 0.9") =====
    # Navy header band
    _add_rect(slide, 0, 0, prs.slide_width, Inches(0.9), COLOR_NAVY)
    # Gold accent bar under header
    _add_rect(slide, 0, Inches(0.9), prs.slide_width, Inches(0.06), COLOR_ACCENT)

    # Left: Topaz wordmark
    _add_text(slide, Inches(0.3), Inches(0.12), Inches(3.5), Inches(0.4),
              "TOPAZ TRANSFORMA", size=16, bold=True, color=COLOR_WHITE)
    _add_text(slide, Inches(0.3), Inches(0.50), Inches(3.5), Inches(0.28),
              "Equity Research — Internal", size=9, italic=True,
              color=COLOR_ACCENT)

    # Center: ticker + company + tagline
    tagline = one_pager_content.get("tagline") or quote.industry or ""
    title = f"{stmt.ticker}   |   {quote.company_name or stmt.title}"
    _add_text(slide, Inches(4.0), Inches(0.12), Inches(6.5), Inches(0.4),
              title, size=16, bold=True, color=COLOR_WHITE,
              align=PP_ALIGN.CENTER)
    if tagline:
        _add_text(slide, Inches(4.0), Inches(0.50), Inches(6.5), Inches(0.28),
                  tagline, size=10, italic=True, color=COLOR_LIGHT,
                  align=PP_ALIGN.CENTER)

    # Right: Recommendation chip
    es = one_pager_content.get("executive_summary") or {}
    recommendation = es.get("recommendation", "REVIEW")
    # Price target: prefer DCF Base; else analyst mean
    dcf_base_price = dcf.scenarios["Base"].price_per_share if dcf and "Base" in dcf.scenarios else None
    dcf_bear_price = dcf.scenarios["Bear"].price_per_share if dcf and "Bear" in dcf.scenarios else None
    dcf_bull_price = dcf.scenarios["Bull"].price_per_share if dcf and "Bull" in dcf.scenarios else None
    chip_target = dcf_base_price or (quote.analyst_target_mean if quote else None)
    upside = (chip_target / quote.current_price - 1) if (chip_target and quote.current_price) else None
    _recommendation_chip(slide, Inches(10.7), Inches(0.1),
                         Inches(2.4), Inches(0.72),
                         recommendation, chip_target, upside)

    # ===== KPI STRIP (1.05 – 1.95") =====
    kpi_y = Inches(1.1)
    kpi_h = Inches(0.82)
    kpi_left = Inches(0.3)
    kpi_gap = Inches(0.12)
    kpi_w = (prs.slide_width - Inches(0.6) - kpi_gap * 4) / 5   # 5 kpis with 4 gaps

    kpis = [
        ("Current Price", f"${quote.current_price:.2f}" if quote.current_price else "—"),
        ("Market Cap", f"${quote.market_cap/1e9:.2f}B" if quote.market_cap else "—"),
        ("Revenue (TTM)", f"${quote.revenue_ttm/1e9:.2f}B" if quote.revenue_ttm else "—"),
        ("EV / EBITDA", f"{quote.ev_ebitda:.1f}x" if quote.ev_ebitda else "—"),
        ("Fwd P/E", f"{quote.forward_pe:.1f}x" if quote.forward_pe else "—"),
    ]
    for i, (label, value) in enumerate(kpis):
        _add_kpi_cell(slide, kpi_left + (kpi_w + kpi_gap) * i, kpi_y,
                      kpi_w, kpi_h, label, value)

    # ===== MAIN BODY (2.1 – 6.8") =====
    body_y = Inches(2.15)

    # ---------- LEFT COLUMN (0.3 - 5.0) — Thesis + Strengths + Risks ----------
    left_x = Inches(0.3)
    left_w = Inches(4.6)
    cursor = body_y

    # Investment Thesis
    intro = one_pager_content.get("intro_paragraph", "")
    if not intro:
        intro = es.get("thesis", f"{stmt.title} ({stmt.ticker}) — see full equity research report for thesis.")
    _add_section_header(slide, left_x, cursor, left_w, "Investment Thesis")
    cursor += Inches(0.34)
    _add_text(slide, left_x, cursor, left_w, Inches(1.15),
              intro, size=9.5, color=COLOR_BLACK)
    cursor += Inches(1.20)

    # Key Strengths
    strengths = (one_pager_content.get("key_strengths") or
                 es.get("key_strengths", []))[:5]
    if strengths:
        _add_section_header(slide, left_x, cursor, left_w, "Key Strengths",
                            color=COLOR_GREEN)
        cursor += Inches(0.34)
        bh = Inches(0.22) * max(len(strengths), 1) + Inches(0.2)
        _add_bullet_block(slide, left_x, cursor, left_w, bh, strengths)
        cursor += bh + Inches(0.05)

    # Key Risks
    risks = (one_pager_content.get("key_risks") or
             es.get("key_risks", []))[:5]
    if risks:
        _add_section_header(slide, left_x, cursor, left_w, "Key Risks",
                            color=COLOR_RED)
        cursor += Inches(0.34)
        bh = Inches(0.22) * max(len(risks), 1) + Inches(0.2)
        _add_bullet_block(slide, left_x, cursor, left_w, bh, risks)

    # ---------- MIDDLE COLUMN (5.1 - 8.7) — Valuation Snapshot + Football ----------
    mid_x = Inches(5.05)
    mid_w = Inches(3.5)

    # Valuation Snapshot box
    val_h = Inches(2.55)
    _valuation_snapshot(slide, mid_x, body_y, mid_w, val_h,
                         current_price=quote.current_price,
                         dcf_bear=dcf_bear_price, dcf_base=dcf_base_price,
                         dcf_bull=dcf_bull_price,
                         analyst_low=quote.analyst_target_low if quote else None,
                         analyst_mean=quote.analyst_target_mean if quote else None,
                         analyst_high=quote.analyst_target_high if quote else None,
                         analyst_count=quote.analyst_count if quote else None)

    # Football field below valuation snapshot
    fb_y = body_y + val_h + Inches(0.12)
    fb_w = mid_w + Inches(0.0)
    fb_h = Inches(2.0)
    if football_field_png and football_field_png.exists():
        _add_text(slide, mid_x, fb_y, fb_w, Inches(0.25),
                  "FOOTBALL FIELD", size=9, bold=True, color=COLOR_NAVY)
        slide.shapes.add_picture(str(football_field_png), mid_x,
                                 fb_y + Inches(0.28), width=fb_w, height=fb_h - Inches(0.28))

    # ---------- RIGHT COLUMN (8.8 - 13.0) — Revenue chart + metrics ----------
    right_x = Inches(8.7)
    right_w = Inches(4.4)

    _add_section_header(slide, right_x, body_y, right_w,
                        "Revenue — Historical + Consensus")
    chart_y = body_y + Inches(0.34)
    chart_h = Inches(2.3)

    rev_series = _quarterly_revenue_series(stmt)
    kpi_series = one_pager_content.get("primary_kpi_series") or None
    if kpi_series and isinstance(kpi_series, dict) and len(kpi_series) == 0:
        kpi_series = None
    kpi_label = one_pager_content.get("primary_kpi_label", "KPI")

    if rev_series and output_path is not None:
        chart_png = output_path.parent / f"{stmt.ticker}_chart.png"
        revenue_kpi_chart(rev_series, kpi_series, kpi_label, chart_png)
        slide.shapes.add_picture(str(chart_png), right_x, chart_y,
                                 width=right_w, height=chart_h)

    # Secondary metrics table under the chart
    metrics_y = chart_y + chart_h + Inches(0.2)
    _add_section_header(slide, right_x, metrics_y, right_w, "Financial Metrics")
    metrics_y += Inches(0.32)

    metrics = []
    if quote.enterprise_value:
        metrics.append(("Enterprise Value", f"${quote.enterprise_value/1e9:.2f}B"))
    if quote.ebitda_ttm:
        metrics.append(("EBITDA (TTM)", f"${quote.ebitda_ttm/1e9:.2f}B"))
    if quote.free_cash_flow:
        metrics.append(("FCF (TTM)", f"${quote.free_cash_flow/1e9:.2f}B"))
    if quote.ev_sales:
        metrics.append(("EV / Sales", f"{quote.ev_sales:.2f}x"))
    if quote.net_leverage is not None:
        metrics.append(("Net Debt / EBITDA", f"{quote.net_leverage:.2f}x"))
    if quote.fcf_yield:
        metrics.append(("FCF Yield", f"{quote.fcf_yield*100:.2f}%"))
    if quote.beta is not None:
        metrics.append(("Beta", f"{quote.beta:.2f}"))
    if quote.week52_low and quote.week52_high:
        metrics.append(("52-Week Range", f"${quote.week52_low:.2f} — ${quote.week52_high:.2f}"))

    # 2-column layout inside the right panel
    half = right_w / 2
    row_h = Inches(0.21)
    mid_idx = (len(metrics) + 1) // 2
    for i, (k, v) in enumerate(metrics[:mid_idx]):
        y = metrics_y + row_h * i
        _add_text(slide, right_x, y, half - Inches(0.05), row_h,
                  k, size=9, color=COLOR_GRAY, vertical=MSO_ANCHOR.MIDDLE)
        _add_text(slide, right_x + half * 0.5, y, half * 0.6, row_h,
                  v, size=9, bold=True, color=COLOR_NAVY_D,
                  align=PP_ALIGN.RIGHT, vertical=MSO_ANCHOR.MIDDLE)
    for i, (k, v) in enumerate(metrics[mid_idx:]):
        y = metrics_y + row_h * i
        _add_text(slide, right_x + half + Inches(0.1), y, half - Inches(0.15), row_h,
                  k, size=9, color=COLOR_GRAY, vertical=MSO_ANCHOR.MIDDLE)
        _add_text(slide, right_x + half * 1.5, y, half * 0.5, row_h,
                  v, size=9, bold=True, color=COLOR_NAVY_D,
                  align=PP_ALIGN.RIGHT, vertical=MSO_ANCHOR.MIDDLE)

    # ===== FOOTER (7.05 - 7.5") =====
    _add_rect(slide, 0, Inches(7.05), prs.slide_width, Inches(0.45), COLOR_NAVY_D)
    footer_txt = (f"Sources: SEC EDGAR • Yahoo Finance (analyst consensus) • "
                  f"Topaz internal DCF   |   "
                  f"Generated {date.today().isoformat()}   |   "
                  f"Confidential — Topaz Transforma Internal Use Only")
    _add_text(slide, Inches(0.3), Inches(7.15), prs.slide_width - Inches(0.6),
              Inches(0.25), footer_txt, size=7.5, italic=True,
              color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
